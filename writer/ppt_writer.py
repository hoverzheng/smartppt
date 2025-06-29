from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

class PPTWriter:
    def __init__(self):
        self.default_layouts = {
            "默认": 1,  # 标题和内容
            "简约": 6,  # 仅标题
            "商务": 1   # 标题和内容
        }
    
    def write_ppt(self, formatted_content, output_path: str, style: str = None):
        """
        根据格式化内容和样式生成PPT文件。
        """
        prs = Presentation()
        layout_index = self.default_layouts.get(style, 1)
        
        for page in formatted_content:
            # 使用标题和内容版式
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # 先创建自定义内容框，再设置标题
            self._create_custom_content_box(slide, page)
            self._set_slide_title(slide, page["title"])
        
        prs.save(output_path)
        return output_path

    def _set_slide_title(self, slide, title_text):
        """设置幻灯片标题"""
        try:
            print(f"正在设置标题: {title_text}")
            
            # 检查是否有标题占位符
            title_placeholder = None
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format'):
                    print(f"找到占位符，类型: {shape.placeholder_format.type}, 名称: {shape.placeholder_format.type}")
                    if shape.placeholder_format.type == 1:  # 标题占位符通常是1
                        title_placeholder = shape
                        break
            
            if title_placeholder:
                print(f"找到标题占位符，设置标题")
                title_placeholder.text = title_text
                self._format_title(title_placeholder)
            else:
                print("没有找到标题占位符，手动创建标题")
                # 手动创建标题文本框
                left = Inches(0.5)
                top = Inches(0.2)
                width = Inches(9.0)
                height = Inches(0.8)
                
                title_box = slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_frame.clear()
                p = title_frame.add_paragraph()
                p.text = title_text
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(44, 62, 80)
                p.alignment = PP_ALIGN.CENTER
        except Exception as e:
            print(f"设置标题时出错: {e}")

    def _format_title(self, title_shape):
        """格式化标题样式"""
        try:
            title_frame = title_shape.text_frame
            p = title_frame.paragraphs[0]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(44, 62, 80)
            p.alignment = PP_ALIGN.CENTER
        except Exception as e:
            print(f"格式化标题时出错: {e}")

    def _create_custom_content_box(self, slide, page_content):
        """删除原有占位符，创建自定义文本框"""
        try:
            print("开始创建自定义内容框")
            
            # 先记录原有的占位符，再创建新的文本框
            original_placeholders = []
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format'):
                    print(f"检查占位符，类型: {shape.placeholder_format.type}")
                    original_placeholders.append(shape)
            
            # 创建新的文本框
            left = Inches(0.5)
            top = Inches(1.2)  # 为标题留出空间
            width = Inches(9.0)
            height = Inches(5.5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            self._fill_content_with_formatting(txBox, page_content)
            
            # 删除原有的内容占位符（保留标题占位符）
            shapes_to_remove = []
            for shape in original_placeholders:
                # 只删除内容占位符，保留标题占位符
                if shape.placeholder_format.type == 1:  # 标题占位符
                    print(f"保留标题占位符")
                else:
                    # 其他类型的占位符（如内容占位符）
                    print(f"找到内容占位符，准备删除")
                    shapes_to_remove.append(shape)
            
            # 删除找到的内容占位符
            for shape in shapes_to_remove:
                print(f"删除内容占位符")
                slide.shapes._spTree.remove(shape._element)
            
        except Exception as e:
            print(f"创建自定义文本框时出错: {e}")

    def _fill_content_with_formatting(self, content_shape, page_content):
        """直接填充格式化内容"""
        try:
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            # 设置文本框属性
            text_frame.word_wrap = True
            
            # 添加总结
            if "summary" in page_content:
                p = text_frame.add_paragraph()
                p.text = f"📋 {page_content['summary']}"
                try:
                    p.font.size = Pt(14)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(52, 73, 94)
                    p.space_after = Pt(8)
                    p.alignment = PP_ALIGN.LEFT
                except Exception as e:
                    print(f"格式化总结段落时出错: {e}")
            
            # 添加要点
            if "points" in page_content:
                print(f"页面包含 {len(page_content['points'])} 个论点")
                for i, point in enumerate(page_content["points"], 1):
                    if isinstance(point, dict) and "main_point" in point:
                        # 新格式：包含主要论点和支持事实
                        p = text_frame.add_paragraph()
                        p.text = f"{i}. {point['main_point']}"
                        try:
                            p.font.size = Pt(13)
                            p.font.bold = True
                            p.font.color.rgb = RGBColor(44, 62, 80)
                            p.space_after = Pt(4)
                            p.space_before = Pt(6)
                            p.alignment = PP_ALIGN.LEFT
                        except Exception as e:
                            print(f"格式化主要论点时出错: {e}")
                        
                        # 添加支持事实
                        if "supporting_facts" in point:
                            facts = point["supporting_facts"]
                            for fact in facts:
                                try:
                                    if isinstance(fact, dict) and "fact" in fact and "explanation" in fact:
                                        # 新格式：包含事实和说明，用冒号连接
                                        run = p.add_run()
                                        run.text = f"\n   • {fact['fact']}: {fact['explanation']}"
                                        run.font.size = Pt(11)
                                        run.font.color.rgb = RGBColor(52, 73, 94)
                                    else:
                                        # 旧格式：简单事实
                                        run = p.add_run()
                                        run.text = f"\n   • {fact}"
                                        run.font.size = Pt(11)
                                        run.font.color.rgb = RGBColor(52, 73, 94)
                                except Exception as e:
                                    print(f"添加支持事实时出错: {e}")
                    else:
                        # 旧格式：简单要点
                        p = text_frame.add_paragraph()
                        p.text = f"{i}. {point}"
                        try:
                            p.font.size = Pt(13)
                            p.font.bold = True
                            p.font.color.rgb = RGBColor(44, 62, 80)
                            p.space_after = Pt(6)
                            p.alignment = PP_ALIGN.LEFT
                        except Exception as e:
                            print(f"格式化简单要点时出错: {e}")
        except Exception as e:
            print(f"填充内容时出错: {e}")

    def write_ppt_with_template(self, formatted_content, template_path: str, output_path: str, style: str = None):
        """
        基于用户上传的模板生成PPT文件。
        """
        try:
            # 加载模板
            prs = Presentation(template_path)
            
            # 分析模板中的版式
            layouts = prs.slide_layouts
            
            # 选择合适的版式
            title_layout = self._find_best_layout(layouts, "title")
            content_layout = self._find_best_layout(layouts, "content")
            
            # 创建新演示文稿
            new_prs = Presentation(template_path)
            
            # 添加内容到幻灯片
            for i, page in enumerate(formatted_content):
                # 选择版式（第一页用标题版式，其他用内容版式）
                if i == 0 and title_layout is not None:
                    slide_layout = title_layout
                else:
                    slide_layout = content_layout or new_prs.slide_layouts[1]
                
                slide = new_prs.slides.add_slide(slide_layout)
                
                # 填充内容
                self._fill_slide_content_with_template(slide, page)
            
            new_prs.save(output_path)
            return output_path
            
        except Exception as e:
            print(f"使用模板生成PPT失败: {e}")
            # 回退到默认方法
            return self.write_ppt(formatted_content, output_path, style)

    def _fill_slide_content_with_template(self, slide, page_content):
        """为模板填充幻灯片内容"""
        try:
            # 先创建自定义内容框，再设置标题
            self._create_custom_content_box(slide, page_content)
            self._set_slide_title(slide, page_content["title"])
            
        except Exception as e:
            print(f"填充模板幻灯片内容时出错: {e}")

    def _find_best_layout(self, layouts, layout_type):
        """
        在模板中查找最适合的版式。
        """
        if layout_type == "title":
            # 查找标题版式（通常只有标题）
            for i, layout in enumerate(layouts):
                if layout.name.lower() in ['title slide', '标题幻灯片', 'title']:
                    return layout
                # 检查占位符数量
                if len(layout.placeholders) == 1:
                    return layout
        else:
            # 查找内容版式（有标题和内容）
            for i, layout in enumerate(layouts):
                if layout.name.lower() in ['title and content', '标题和内容', 'content']:
                    return layout
                # 检查占位符数量
                if len(layout.placeholders) >= 2:
                    return layout
        
        # 如果没找到合适的，返回第一个版式
        return layouts[0] if layouts else None

    def get_template_info(self, template_path: str):
        """
        获取模板信息。
        """
        try:
            prs = Presentation(template_path)
            layouts = prs.slide_layouts
            
            layout_info = []
            for i, layout in enumerate(layouts):
                layout_info.append({
                    "index": i,
                    "name": layout.name,
                    "placeholders": len(layout.placeholders)
                })
            
            return {
                "total_layouts": len(layouts),
                "layouts": layout_info
            }
        except Exception as e:
            return {"error": str(e)} 